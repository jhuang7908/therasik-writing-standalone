#!/usr/bin/env python3
"""Reference ContentDoc -> editable .pptx renderer (Creative Studio MVP).

Implements PPTX_TEMPLATE_SPEC.md. Text is written into native, editable
text boxes (never baked into an image). Images (visual.resolved) are
optional; missing visuals fall back to palette color blocks.

Run with the built-in default template (no external files needed)::

    python renderer/ppt_renderer.py --doc ../schemas/example_ppt.json \
        --out ../outputs/example_ppt.pptx

Dependencies: python-pptx (already in the WSL venv_linux).
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from pptx.enum.dml import MSO_THEME_COLOR
import lxml.etree as etree


# --------------------------------------------------------------------------- #
# Built-in default template (matches PPTX_TEMPLATE_SPEC.md, ppt_tech_blue_01)
# --------------------------------------------------------------------------- #
DEFAULT_TEMPLATE = {
    "template_id": "ppt_tech_blue_01",
    "slide_size": {"w_in": 13.333, "h_in": 7.5},
    "palette": {
        "primary":    "#1E3A8A",   # deep navy
        "secondary":  "#0EA5E9",   # sky blue
        "accent":     "#F59E0B",   # amber
        "text_dark":  "#1F2937",
        "text_light": "#FFFFFF",
        "bg":         "#F0F6FF",   # very light blue-white
        "bg_dark":    "#0F172A",   # near-black
        "muted":      "#94A3B8",
    },
    "gradient": {
        "cover":   ["#1E3A8A", "#0F172A"],   # dark navy gradient
        "section": ["#1E3A8A", "#1E40AF"],
        "content": ["#F0F6FF", "#E0ECFF"],
    },
    "fonts": {"title": "思源黑体", "body": "思源黑体", "fallback": "Microsoft YaHei"},
    "layouts": {
        "cover": {
            "bg_gradient": ["primary", "bg_dark"],
            "accent_bar": {"vertical": True, "x": 0.0, "w": 0.07, "color": "secondary"},
            "bottom_bar": True,
            "boxes": {
                "tag":      {"x": 1.0,  "y": 1.5,  "w": 9.0,  "h": 0.45, "size": 14,
                             "color": "secondary", "align": "left"},
                "title":    {"x": 1.0,  "y": 2.1,  "w": 10.8, "h": 2.4,  "size": 46,
                             "bold": True, "color": "text_light", "align": "left"},
                "subtitle": {"x": 1.0,  "y": 4.7,  "w": 10.0, "h": 0.9,  "size": 22,
                             "color": "secondary", "align": "left"},
            },
        },
        "section": {
            "bg_gradient": ["primary", "secondary"],
            "accent_bar": {"vertical": False, "y": 0.0, "h": 0.08, "color": "accent"},
            "bottom_bar": True,
            "boxes": {
                "title":    {"x": 1.0,  "y": 2.6,  "w": 11.3, "h": 1.6,  "size": 40,
                             "bold": True, "color": "text_light", "align": "left"},
                "subtitle": {"x": 1.0,  "y": 4.4,  "w": 11.3, "h": 0.9,  "size": 22,
                             "color": "text_light", "align": "left"},
            },
        },
        "title_bullets": {
            "bg_gradient": ["bg", "bg"],
            "accent_bar": {"vertical": True, "x": 0.0, "w": 0.055, "color": "primary"},
            "bottom_bar": True,
            "boxes": {
                "title":   {"x": 0.85, "y": 0.55, "w": 11.7, "h": 1.05, "size": 34,
                            "bold": True, "color": "primary", "align": "left"},
                "bullets": {"x": 1.05, "y": 1.85, "w": 11.2, "h": 4.85, "size": 20,
                            "color": "text_dark", "line_spacing": 1.6},
            },
        },
        "two_column": {
            "bg_gradient": ["bg", "bg"],
            "accent_bar": {"vertical": True, "x": 0.0, "w": 0.055, "color": "secondary"},
            "bottom_bar": True,
            "boxes": {
                "title":   {"x": 0.85, "y": 0.55, "w": 11.7, "h": 1.0,  "size": 34,
                            "bold": True, "color": "primary", "align": "left"},
                "body":    {"x": 0.85, "y": 1.8,  "w": 5.8,  "h": 4.8,  "size": 18,
                            "color": "text_dark"},
                "bullets": {"x": 7.0,  "y": 1.8,  "w": 5.8,  "h": 4.8,  "size": 18,
                            "color": "text_dark", "line_spacing": 1.5},
            },
        },
        "image_left": {
            "bg_gradient": ["bg", "bg"],
            "image_region": {"x": 0.0, "y": 0.0, "w": 6.2, "h": 7.5},
            "bottom_bar": True,
            "boxes": {
                "title":   {"x": 6.8,  "y": 1.2,  "w": 6.0,  "h": 1.1,  "size": 30,
                            "bold": True, "color": "primary"},
                "body":    {"x": 6.8,  "y": 2.5,  "w": 6.0,  "h": 3.8,  "size": 18,
                            "color": "text_dark"},
                "caption": {"x": 6.8,  "y": 6.5,  "w": 6.0,  "h": 0.45, "size": 12,
                            "color": "muted"},
            },
        },
        "image_right": {
            "bg_gradient": ["bg", "bg"],
            "image_region": {"x": 7.133, "y": 0.0, "w": 6.2, "h": 7.5},
            "accent_bar": {"vertical": True, "x": 0.0, "w": 0.055, "color": "primary"},
            "bottom_bar": True,
            "boxes": {
                "title":   {"x": 0.7,  "y": 1.2,  "w": 6.0,  "h": 1.1,  "size": 30,
                            "bold": True, "color": "primary"},
                "body":    {"x": 0.7,  "y": 2.5,  "w": 6.0,  "h": 3.8,  "size": 18,
                            "color": "text_dark"},
                "caption": {"x": 0.7,  "y": 6.5,  "w": 6.0,  "h": 0.45, "size": 12,
                            "color": "muted"},
            },
        },
        "big_number": {
            "bg_gradient": ["bg_dark", "primary"],
            "bottom_bar": True,
            "boxes": {
                "title":    {"x": 0.9, "y": 0.65, "w": 11.5, "h": 1.0,  "size": 30,
                             "bold": True, "color": "text_light"},
                "subtitle": {"x": 0.9, "y": 1.7,  "w": 11.5, "h": 0.65, "size": 18,
                             "color": "secondary"},
                "bullets":  {"x": 0.9, "y": 2.7,  "w": 11.5, "h": 2.8,  "size": 52,
                             "bold": True, "color": "accent", "line_spacing": 1.2},
                "caption":  {"x": 0.9, "y": 6.3,  "w": 11.5, "h": 0.5,  "size": 14,
                             "color": "muted"},
            },
        },
        "quote": {
            "bg_gradient": ["primary", "bg_dark"],
            "accent_bar": {"vertical": False, "y": 0.0, "h": 0.08, "color": "accent"},
            "bottom_bar": True,
            "boxes": {
                "body":    {"x": 1.4, "y": 2.3,  "w": 10.5, "h": 2.6,  "size": 32,
                            "bold": True, "color": "text_light", "align": "center"},
                "caption": {"x": 1.4, "y": 5.1,  "w": 10.5, "h": 0.65, "size": 18,
                            "color": "secondary", "align": "center"},
            },
        },
    },
}


def _hex(color: str, palette: dict) -> RGBColor:
    val = palette.get(color, color)
    return RGBColor.from_string(val.lstrip("#"))


def _align(name: str) -> PP_ALIGN:
    return {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT}.get(name, PP_ALIGN.LEFT)


def _set_cn_font(run, font_name: str) -> None:
    """Set both Latin and East-Asian font so CJK doesn't fall back to tofu."""
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font_name)


def _add_bg(slide, prs, color: RGBColor) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def _add_gradient_bg(slide, prs, color1: RGBColor, color2: RGBColor) -> None:
    """Full-slide linear gradient background (top-left → bottom-right)."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 135.0   # diagonal
    stops = fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = color1
    stops[1].position = 1.0
    stops[1].color.rgb = color2
    shape.line.fill.background()
    shape.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def _add_accent_bar(slide, prs, color: RGBColor,
                    x=0.0, y=0.0, w=0.06, h=None, vertical=True) -> None:
    """Thin colored bar — left edge accent or top strip."""
    from pptx.enum.shapes import MSO_SHAPE
    sw = prs.slide_width;  sh = prs.slide_height
    if vertical:
        bar_w = Inches(w); bar_h = sh
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), 0, bar_w, bar_h)
    else:
        bar_w = sw; bar_h = Inches(h or 0.05)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), bar_w, bar_h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background(); shape.shadow.inherit = False


def _add_bottom_bar(slide, prs, bg_color: RGBColor, accent: RGBColor,
                    title_text: str = "", page_num: int = 0) -> None:
    """Footer strip with optional title echo and page number."""
    from pptx.enum.shapes import MSO_SHAPE
    sh = prs.slide_height; sw = prs.slide_width; bh = Inches(0.36)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, sh - bh, sw, bh)
    bar.fill.solid(); bar.fill.fore_color.rgb = bg_color
    bar.line.fill.background(); bar.shadow.inherit = False
    # accent dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.25), sh - bh + Inches(0.1),
                                  Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = accent
    dot.line.fill.background()
    if page_num:
        tb = slide.shapes.add_textbox(sw - Inches(0.6), sh - bh, Inches(0.5), bh)
        tf = tb.text_frame; para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.RIGHT
        run = para.add_run(); run.text = str(page_num)
        run.font.size = Pt(10); run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.bold = False


def _add_image_to_slide(slide, image_path: str,
                         x: float, y: float, w: float, h: float) -> None:
    """Insert image from local path, fail silently if missing."""
    try:
        p = Path(image_path)
        if p.exists():
            slide.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(w), Inches(h))
    except Exception as exc:
        logger.warning("Image insert failed (%s): %s", image_path, exc)


def _add_text(slide, box: dict, palette: dict, fonts: dict,
              text: str = "", bullets: Optional[list] = None) -> None:
    tb = slide.shapes.add_textbox(
        Inches(box["x"]), Inches(box["y"]),
        Inches(box["w"]), Inches(box["h"]))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    color = _hex(box.get("color", "text_dark"), palette)
    size = Pt(box.get("size", 18))
    bold = box.get("bold", False)
    align = _align(box.get("align", "left"))
    font_name = fonts.get("body", "思源黑体")
    line_spacing = box.get("line_spacing")

    items = bullets if bullets else [text]
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if line_spacing:
            para.line_spacing = line_spacing
        prefix = "• " if bullets else ""
        run = para.add_run()
        run.text = f"{prefix}{item}"
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
        _set_cn_font(run, font_name)


def _normalize_block(block: dict) -> dict:
    """Map ContentDoc field names (headline/subheadline/body_paragraphs)
    to renderer field names (title/subtitle/bullets).
    Also handles layout aliases from LLM output.
    """
    LAYOUT_ALIASES = {
        "title_hero":    "cover",
        "title_slide":   "cover",
        "hero":          "cover",
        "cover_slide":   "cover",
        "section_break": "section",
        "divider":       "section",
        "content":       "title_bullets",
        "bullets":       "title_bullets",
        "bullet_points": "title_bullets",
        "text_only":     "title_bullets",
        "split":         "two_column",
        "two_col":       "two_column",
        "split_layout":  "two_column",
        "image_left_text":"image_left",
        "image_right_text":"image_right",
        "stats":         "big_number",
        "numbers":       "big_number",
        "closing":       "cover",
        "summary":       "title_bullets",
        "thank_you":     "cover",
    }

    block = dict(block)
    # Normalize layout name
    layout = block.get("layout", "title_bullets")
    block["layout"] = LAYOUT_ALIASES.get(layout, layout)

    # Normalize text fields
    text = dict(block.get("text", {}) or {})
    # headline → title
    if "headline" in text and "title" not in text:
        text["title"] = text.pop("headline")
    # subheadline → subtitle
    if "subheadline" in text and "subtitle" not in text:
        text["subtitle"] = text.pop("subheadline")
    # body_paragraphs → bullets (list of strings)
    if "body_paragraphs" in text and not block.get("bullets"):
        paras = text.pop("body_paragraphs")
        if isinstance(paras, list):
            block["bullets"] = [str(p) for p in paras if p]
        elif isinstance(paras, str) and paras:
            block["bullets"] = [paras]
    # body string → bullets
    if "body" in text and not block.get("bullets"):
        body = text.get("body", "")
        if isinstance(body, str) and body:
            block["bullets"] = [body]
    # caption → subtitle fallback
    if "caption" in text and "subtitle" not in text:
        text["subtitle"] = text.get("caption", "")

    block["text"] = text
    return block


def render_ppt(doc: dict, template: Optional[dict], out_path: str) -> dict:
    template = template or DEFAULT_TEMPLATE
    palette = template["palette"]
    fonts = template["fonts"]
    layouts = template["layouts"]

    prs = Presentation()
    prs.slide_width = Inches(template["slide_size"]["w_in"])
    prs.slide_height = Inches(template["slide_size"]["h_in"])
    blank = prs.slide_layouts[6]

    qa = {"slide_count": 0, "missing_layout": [], "empty_title": [],
          "font_set": True, "failures": []}

    for raw_block in doc["blocks"]:
        block = _normalize_block(raw_block)
        slide_num = qa["slide_count"] + 1
        layout_name = block["layout"]
        spec = layouts.get(layout_name)
        if spec is None:
            qa["missing_layout"].append(layout_name)
            spec = layouts.get("title_bullets")

        slide = prs.slides.add_slide(blank)

        # ── Background ──────────────────────────────────────────────────────
        grad = spec.get("bg_gradient")
        if grad and len(grad) == 2:
            c1 = _hex(grad[0], palette)
            c2 = _hex(grad[1], palette)
            if c1 == c2:
                _add_bg(slide, prs, c1)
            else:
                _add_gradient_bg(slide, prs, c1, c2)
        else:
            _add_bg(slide, prs, _hex(spec.get("bg_fill", "bg"), palette))

        # ── Accent bar ───────────────────────────────────────────────────────
        ab = spec.get("accent_bar")
        if ab:
            _add_accent_bar(slide, prs, _hex(ab["color"], palette),
                            x=ab.get("x", 0), y=ab.get("y", 0),
                            w=ab.get("w", 0.06), h=ab.get("h", 0.06),
                            vertical=ab.get("vertical", True))

        # ── Image region ─────────────────────────────────────────────────────
        img_region = spec.get("image_region")
        if img_region:
            vis = block.get("visual") or {}
            resolved = vis.get("resolved") or {}
            img_path = resolved.get("url")
            if img_path:
                _add_image_to_slide(slide, img_path,
                                     img_region["x"], img_region["y"],
                                     img_region["w"], img_region["h"])
            else:
                # placeholder color block
                from pptx.enum.shapes import MSO_SHAPE
                ph = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(img_region["x"]), Inches(img_region["y"]),
                    Inches(img_region["w"]), Inches(img_region["h"]))
                ph.fill.solid()
                ph.fill.fore_color.rgb = _hex("primary", palette)
                ph.fill.fore_color.theme_color  # touch to avoid rounding
                ph.line.fill.background()

        # ── Text boxes ───────────────────────────────────────────────────────
        text = block.get("text", {}) or {}
        bullets = block.get("bullets", []) or []
        boxes = spec.get("boxes", {})

        for key, box in boxes.items():
            if key == "bullets":
                if bullets:
                    _add_text(slide, box, palette, fonts, bullets=bullets)
            else:
                val = text.get(key, "")
                if val:
                    _add_text(slide, box, palette, fonts, text=val)

        # ── Bottom bar ───────────────────────────────────────────────────────
        if spec.get("bottom_bar"):
            _add_bottom_bar(slide, prs,
                             _hex("primary", palette),
                             _hex("secondary", palette),
                             page_num=slide_num)

        if not text.get("title") and layout_name not in ("quote", "section"):
            qa["empty_title"].append(block.get("block_id", f"slide_{slide_num}"))

        qa["slide_count"] += 1

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)

    qa["expected"] = len(doc["blocks"])
    qa["ok"] = (qa["slide_count"] == qa["expected"]
                and not qa["empty_title"])
    qa["out"] = out_path
    return qa


def main() -> int:
    p = argparse.ArgumentParser(description="ContentDoc -> editable PPTX")
    p.add_argument("--doc", required=True, help="ContentDoc JSON path")
    p.add_argument("--template", help="template.json path (optional)")
    p.add_argument("--out", required=True, help="output .pptx path")
    args = p.parse_args()

    doc = json.loads(Path(args.doc).read_text(encoding="utf-8"))
    template = (json.loads(Path(args.template).read_text(encoding="utf-8"))
                if args.template else None)
    qa = render_ppt(doc, template, args.out)
    print(json.dumps(qa, ensure_ascii=False, indent=2))
    return 0 if qa["ok"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
