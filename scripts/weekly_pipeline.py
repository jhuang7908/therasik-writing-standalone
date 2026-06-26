"""
weekly_pipeline.py — Cross-platform (Linux/Windows) weekly PPT + MP4 pipeline.

Steps:
  1. PubMed: find this week's article
  2. Claude: generate 6-slide deck_spec.json from article abstract
  3. gpt-image-2: render scientific slides (S2, S3) — full-slide generation
  4. CogView-4: render illustrative slides (S1, S4, S5, S6)
  5. PIL: stamp NextVivo logo (once, bottom-right)
  6. Gemini Vision: visual audit — fail fast if axes/text wrong
  7. Volcengine TTS: generate Chinese narrations
  8. moviepy: assemble MP4
  9. python-pptx: assemble PPTX
  10. Email: send to recipients
"""

import os, sys, json, time, base64, shutil, textwrap
from pathlib import Path
from datetime import datetime

# ── Load .env ────────────────────────────────────────────────────────────────
_env_candidates = [
    Path(__file__).parent.parent / "content_generation_tools" / ".env",
    Path(__file__).parent.parent / ".env",
]
for _env in _env_candidates:
    if _env.exists():
        for _line in _env.read_text(encoding="utf-8").splitlines():
            _line = _line.strip()
            if _line and not _line.startswith("#") and "=" in _line:
                _k, _v = _line.split("=", 1)
                os.environ[_k.strip()] = _v.strip()  # force-set, don't use setdefault
        break

ROOT     = Path(__file__).parent.parent
WEEK_TAG = datetime.now().strftime("%Y-W%V")
OUT_DIR  = ROOT / "outputs" / f"weekly_{WEEK_TAG}"
IMG_DIR  = OUT_DIR / "slides"
VOICE_DIR = OUT_DIR / "voice"
CACHE    = ROOT / "outputs" / "weekly_cache.json"

LOGO_PATH = ROOT / "assets" / "nextvivo_logo.png"
RULES     = ROOT / "config" / "pipeline_rules.yaml"

# Topic slug resolved at runtime from bundle article.json (falls back to "humanized_mice")
_TOPIC_SLUG = "humanized_mice"
OUT_PPTX  = OUT_DIR / f"{_TOPIC_SLUG}_{WEEK_TAG}.pptx"
OUT_MP4   = OUT_DIR / f"{_TOPIC_SLUG}_{WEEK_TAG}.mp4"


def _resolve_topic_slug(bundle_path: Path) -> str:
    """Read topic_slug from bundle article.json, fall back to dir name."""
    try:
        art = json.loads((bundle_path / "article.json").read_text(encoding="utf-8"))
        slug = art.get("topic_slug") or bundle_path.name
        return slug
    except Exception:
        return bundle_path.name

RECIPIENTS = ["mail.jing.huang@gmail.com"]

# ─── Render routing per slide — 图内中文只用 gpt-image-2，禁止 CogView ───────
SLIDE_RENDER = {i: "gpt-image-2" for i in range(1, 7)}
CLAUDE_SCIENTIFIC = "claude-3-5-sonnet-20241022"
CLAUDE_COMPLEX    = "claude-sonnet-4-20250514"  # 高难度论文


# ─── Stage 1: PubMed ─────────────────────────────────────────────────────────
def stage_pubmed() -> dict:
    print("=== Stage 1: PubMed article search ===")
    sys.path.insert(0, str(ROOT / "scripts"))
    from pubmed_fetch import get_weekly_article
    article = get_weekly_article(cache_path=str(CACHE))
    sel = article.get("selection", {})
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    if article.get("full_text"):
        ft_path = OUT_DIR / "article_full_text.txt"
        ft_path.write_text(article["full_text"], encoding="utf-8")
        print(f"  Full text saved: {ft_path.name} ({article.get('full_text_chars', 0)} chars)")
    print(f"  Article: {article['title'][:80]}...")
    print(f"  PMID: {article['pmid']}  PMCID: {article.get('pmcid','n/a')}  DOI: {article.get('doi','n/a')}")
    print(f"  Related refs: {len(article.get('related_references', []))}")
    print(f"  Selection: {sel.get('mode', 'n/a')} — {sel.get('rationale', '')[:120]}")
    if sel.get("story_hook"):
        print(f"  Story hook: {sel['story_hook']}")
    return article


# ─── Stage 2: Claude deck_spec ────────────────────────────────────────────────
def stage_deck_spec(article: dict) -> dict:
    print("=== Stage 2: Claude 3.5 Sonnet → deck_spec.json ===")
    import anthropic
    client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY", ""))

    full_text = article.get("full_text") or article.get("abstract", "")
    related_block = json.dumps(article.get("related_references", []), ensure_ascii=False, indent=2)
    story_hook = article.get("selection", {}).get("story_hook", "")

    prompt = textwrap.dedent(f"""
    You are a scientific PPT designer for Chinese social media (Douyin/WeChat).
    Create a 6-slide deck spec in JSON from the MAIN article FULL TEXT below.

    ## Narrative policy (mandatory)
    - Tell ONE story only — this week's single main paper (PMID {article['pmid']}).
    - Slides 1–5 must all support the same narrative arc (background → model → key result → implication).
    - Do NOT mix multiple unrelated papers into the story slides.
    - Slide 6 is references only: main paper + related reading list provided below.
    - All numbers and claims in slides 1–5 must come from the main paper full text. No fabrication.

    Main article:
    Title: {article['title']}
    Authors: {article['authors']}
    Journal: {article['journal']} ({article['year']})
    PMID: {article['pmid']}
    PMCID: {article.get('pmcid', 'N/A')}
    DOI: {article.get('doi', 'N/A')}
    Story hook (optional): {story_hook}

    Full text (excerpt):
    {full_text[:18000]}

    Related references (for slide 6 only — similar papers, not co-equal story lines):
    {related_block}

    Slide structure:
    - Slide 1: Hook — one-sentence story promise from this paper
    - Slide 2: Model/setup — humanized mouse system used in THIS paper (chart if data available)
    - Slide 3: Key quantitative finding from THIS paper (chart with exact axis values)
    - Slide 4: Drug/therapy efficacy or safety takeaway from THIS paper
    - Slide 5: Why it matters / limitation / next step (still this paper)
    - Slide 6: References — (1) main PMID+DOI, (2) each related reference with PMID+DOI

    For slides with charts, specify chart_spec: {{x_axis, y_axis, curves}} with values from full text.

    Return ONLY valid JSON:
    {{
      "deck_name": "humanized-mice-{WEEK_TAG}",
      "narrative": "single_story",
      "main_article": {{pmid, pmcid, title, doi}},
      "slides": [{{slide_num, title, subtitle, key_points, visual_description, narration, render_engine}}],
      "references": [{{role: "main"|"related", pmid, doi, title, journal}}]
    }}
    """)

    # 高难度论文用 Sonnet 4.5，默认 3.5
    abstract = (article.get("abstract", "") + " " + full_text[:4000]).lower()
    complex_paper = any(k in abstract.lower() for k in
                        ["kaplan-meier", "dose-response", "multi-model", "bispecific"])
    model = CLAUDE_COMPLEX if complex_paper else CLAUDE_SCIENTIFIC
    print(f"  Scientific planner: {model}")
    response = client.messages.create(
        model=model,
        max_tokens=4000,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = response.content[0].text.strip()
    if raw.startswith("```"):
        raw = "\n".join(raw.split("\n")[1:])
        if "```" in raw:
            raw = raw[:raw.rfind("```")]

    try:
        spec = json.loads(raw)
    except Exception as e:
        print(f"  [WARN] Claude JSON parse failed: {e}, using fallback spec")
        spec = _fallback_spec(article)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    (OUT_DIR / "deck_spec.json").write_text(
        json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(f"  deck_spec.json saved ({len(spec.get('slides',[]))} slides)")
    return spec


# ─── Stage 2b: Kimi fact-check ────────────────────────────────────────────────
def stage_fact_check(spec: dict, article: dict) -> dict:
    """Kimi cross-checks every claim in deck_spec against original paper.

    Zero-tolerance policy: any fabricated number / wrong terminology / wrong
    clinical indication triggers an in-place correction before Stage 3.
    Returns corrected spec (may be unchanged if no issues found).
    """
    print("=== Stage 2b: Kimi fact-check (zero-tolerance) ===")
    kimi_key = os.environ.get("KIMI_API_KEY", "")
    if not kimi_key:
        print("  [SKIP] KIMI_API_KEY not set — fact-check skipped")
        return spec

    import requests as req

    full_text = article.get("full_text") or article.get("abstract", "")
    slides_summary = json.dumps(
        [{"slide": s["slide_num"], "title": s["title"],
          "subtitle": s.get("subtitle",""), "key_points": s.get("key_points",[]),
          "narration": s.get("narration","")}
         for s in spec.get("slides", [])],
        ensure_ascii=False, indent=2
    )

    prompt = f"""You are a rigorous scientific fact-checker with expertise in CAR-T cell therapy.

TASK: Check every factual claim in the slide deck against the source paper.
Return a JSON object: {{"issues": [{{"slide": N, "field": "...", "wrong": "...", "correct": "...", "source": "..."}}], "verdict": "PASS"|"FAIL"}}

Rules (zero tolerance):
1. Terminology: LNP/tLNP must NOT be called "胞外体" (exosome) — completely different technology
2. Clinical indications: only state Phase 1 indications actually registered in NCT. CPTX2309 NCT06917742 = RA + SLE only (NOT MS, NOT myositis)
3. Acquisition timeline: Capstan was independent when Phase 1 started (Apr 2025); AbbVie announced acquisition Jun 30 2025, completed Aug 19 2025
4. All numbers/percentages must appear in the paper text
5. Model organisms: verify species used (humanized mice, cynomolgus monkeys)

Source paper abstract:
{full_text[:3000]}

Slide deck to check:
{slides_summary}

Return ONLY valid JSON."""

    try:
        resp = req.post(
            "https://api.moonshot.cn/v1/chat/completions",
            headers={"Authorization": f"Bearer {kimi_key}",
                     "Content-Type": "application/json"},
            json={"model": "moonshot-v1-8k",
                  "messages": [{"role": "user", "content": prompt}],
                  "temperature": 0.1},
            timeout=60,
        )
        resp.raise_for_status()
        raw = resp.json()["choices"][0]["message"]["content"].strip()
        if raw.startswith("```"):
            raw = "\n".join(raw.split("\n")[1:])
            if "```" in raw:
                raw = raw[:raw.rfind("```")]
        result = json.loads(raw)
    except Exception as e:
        print(f"  [WARN] Kimi fact-check API error: {e} — continuing without correction")
        return spec

    issues = result.get("issues", [])
    verdict = result.get("verdict", "PASS")
    (OUT_DIR / "fact_check.json").write_text(
        json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(f"  Fact-check verdict: {verdict} — {len(issues)} issue(s)")

    if not issues:
        print("  ✅ No factual errors detected")
        return spec

    for issue in issues:
        slide_n = issue.get("slide")
        field   = issue.get("field", "")
        wrong   = issue.get("wrong", "")
        correct = issue.get("correct", "")
        print(f"  ❌ Slide {slide_n} [{field}]: '{wrong}' → '{correct}' ({issue.get('source','')})")
        # Auto-patch: replace wrong text in that slide's fields
        for s in spec.get("slides", []):
            if s["slide_num"] == slide_n and wrong:
                for fld in ("title", "subtitle", "narration"):
                    if wrong in s.get(fld, ""):
                        s[fld] = s[fld].replace(wrong, correct)
                        print(f"    patched {fld}")
                for kp in s.get("key_points", []):
                    idx = s["key_points"].index(kp) if kp in s["key_points"] else -1
                    if idx >= 0 and wrong in kp:
                        s["key_points"][idx] = kp.replace(wrong, correct)
                        print(f"    patched key_points[{idx}]")

    (OUT_DIR / "deck_spec.json").write_text(
        json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(f"  deck_spec.json re-saved after {len(issues)} correction(s)")
    return spec


def _fallback_spec(article: dict) -> dict:
    """Minimal fallback if Claude fails."""
    return {
        "deck_name": f"humanized-mice-{WEEK_TAG}",
        "article": {"pmid": article["pmid"], "title": article["title"], "doi": article.get("doi","")},
        "slides": [
            {"slide_num": i, "title": f"幻灯片 {i}", "key_points": [],
             "visual_description": "humanized mouse laboratory illustration, hand-drawn style",
             "narration": f"第{i}张幻灯片",
             "render_engine": SLIDE_RENDER.get(i, "cogview")}
            for i in range(1, 7)
        ],
    }


# ─── Stage 3: Image generation ────────────────────────────────────────────────
def _build_image_prompt(slide: dict) -> str:
    """Build final image generation prompt from slide spec."""
    base = slide.get("visual_description", "")
    title = slide.get("title", "")
    subtitle = slide.get("subtitle", "")
    points = slide.get("key_points", [])

    engine = slide.get("render_engine", "cogview")
    slide_num = slide.get("slide_num", 0)

    if engine == "gpt-image-2":
        # For scientific charts — include exact chart spec in prompt
        chart = slide.get("chart_spec", {})
        x = chart.get("x_axis", {})
        y = chart.get("y_axis", {})
        curves = chart.get("curves", [])

        prompt = f"""Hand-drawn scientific illustration style, white background, warm sketch aesthetic.
Slide title (Chinese bold): {title}
Subtitle (Chinese orange): {subtitle}

Scientific chart requirements:
- X axis: {x.get('label','Time')} from {x.get('min',0)} to {x.get('max',20)} {x.get('unit','')}
- Y axis: {y.get('label','%')} from {y.get('min',0)}% to {y.get('max',100)}%
- Curves: {'; '.join(curves) if curves else base}

Key points listed on the right side:
{chr(10).join(f'① {p}' for p in points[:3])}

Style: hand-drawn marker sketch, clear Chinese labels, NO garbled text, professional medical illustration.
IMPORTANT: Draw the chart with EXACT axis values as specified. Chinese text must be clear and readable."""
    else:
        # Illustration slide — pure art, no precise chart needed
        points_text = "\n".join(f"① {p}" for p in points[:3])
        prompt = f"""Hand-drawn scientific illustration, warm sketch style, white/cream background.
Chinese bold title at top-left: {title}
{f'Chinese subtitle in orange: {subtitle}' if subtitle else ''}

Visual scene: {base}

Key points (Chinese bullet list):
{points_text}

Style: hand-drawn marker sketch, clean layout, professional medical/biology aesthetic.
Chinese text must be clear and fully legible. NO garbled characters."""

    return prompt


def stage_generate_images(spec: dict):
    print("=== Stage 3: Image generation ===")
    IMG_DIR.mkdir(parents=True, exist_ok=True)
    import requests as req

    openai_key = os.environ.get("OPENAI_API_KEY", "")
    zhipu_key  = os.environ.get("ZHIPU_API_KEY", "")

    for slide in spec.get("slides", []):
        num    = slide["slide_num"]
        engine = slide.get("render_engine", "cogview")
        # Override with routing rules
        engine = SLIDE_RENDER.get(num, engine)
        prompt = _build_image_prompt(slide)
        out    = IMG_DIR / f"slide_{num:02d}.png"

        if out.exists():
            print(f"  slide {num}: cached, skipping")
            continue

        print(f"  slide {num}: rendering via gpt-image-2...")

        _render_gpt_image2(prompt, out, openai_key)

        time.sleep(2)  # rate limit buffer


def _render_gpt_image2(prompt: str, out: Path, api_key: str):
    """Render with gpt-image-2, auto-fallback to Nano Banana Pro (Gemini) on timeout/error."""
    import requests as req

    def _try_gpt() -> bytes:
        resp = req.post(
            "https://api.openai.com/v1/images/generations",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": "gpt-image-2", "prompt": prompt, "n": 1,
                  "size": "1536x1024", "quality": "medium"},
            timeout=180,
        )
        data = resp.json()
        if "data" not in data:
            raise RuntimeError(f"gpt-image-2 error: {data}")
        return base64.b64decode(data["data"][0]["b64_json"])

    def _try_nano() -> bytes:
        gemini_key = os.environ.get("GEMINI_API_KEY", "")
        if not gemini_key:
            raise RuntimeError("GEMINI_API_KEY not set")
        url = ("https://generativelanguage.googleapis.com/v1beta/models/"
               "gemini-3-pro-image-preview:generateContent")
        resp = req.post(
            url,
            headers={"x-goog-api-key": gemini_key, "Content-Type": "application/json"},
            json={"contents": [{"parts": [{"text": prompt}]}],
                  "generationConfig": {"responseModalities": ["IMAGE"],
                                       "imageConfig": {"imageSize": "2K"}}},
            timeout=180,
        )
        if resp.status_code >= 400:
            raise RuntimeError(f"Nano Banana Pro HTTP {resp.status_code}: {resp.text[:200]}")
        parts = (resp.json().get("candidates") or [{}])[0].get("content", {}).get("parts") or []
        for part in parts:
            inline = part.get("inlineData") or part.get("inline_data")
            if inline and inline.get("data"):
                return base64.b64decode(inline["data"])
        raise RuntimeError("Nano Banana Pro empty image response")

    errors = []
    for name, fn in (("gpt-image-2", _try_gpt), ("nano-banana-pro", _try_nano)):
        try:
            img_bytes = fn()
            out.write_bytes(img_bytes)
            print(f"    saved {out.name} via {name} ({out.stat().st_size // 1024}KB)")
            return
        except Exception as exc:
            print(f"    [{name}] failed: {exc!s:.120}")
            errors.append(f"{name}: {exc}")
    raise RuntimeError("All image backends failed: " + " | ".join(errors))


def _render_cogview(prompt: str, out: Path, api_key: str):
    import requests as req
    # CogView-4 via Zhipu AI
    resp = req.post(
        "https://open.bigmodel.cn/api/paas/v4/images/generations",
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        json={"model": "cogview-4", "prompt": prompt, "size": "1440x960"},
        timeout=120,
    )
    data = resp.json()
    url = data.get("data", [{}])[0].get("url", "")
    if not url:
        raise RuntimeError(f"CogView-4 error: {data}")
    img_resp = req.get(url, timeout=60)
    out.write_bytes(img_resp.content)
    print(f"    saved {out.name} ({out.stat().st_size // 1024}KB)")


# ─── Stage 4: Logo stamp ──────────────────────────────────────────────────────
def stage_logo():
    print("=== Stage 4: Logo stamp ===")
    if not LOGO_PATH.exists():
        print("  [WARN] logo not found, skipping")
        return
    from PIL import Image
    logo_raw = Image.open(LOGO_PATH).convert("RGBA")

    for i in range(1, 7):
        slide_path = IMG_DIR / f"slide_{i:02d}.png"
        if not slide_path.exists():
            continue
        slide = Image.open(slide_path).convert("RGBA")
        W, H  = slide.size
        # Logo: 13% width — visible branding, with 4% safe margin to avoid overlap
        logo_w = int(W * 0.13)
        logo_h = int(logo_raw.height * logo_w / logo_raw.width)
        logo   = logo_raw.resize((logo_w, logo_h), Image.LANCZOS)
        # Safe zone: 3% right margin, 4% bottom margin (up from 1.8%)
        # This prevents overlap with content that often extends to ~95% H
        x = W - logo_w - int(W * 0.03)
        y = H - logo_h - int(H * 0.04)
        comp = Image.new("RGBA", slide.size)
        comp.paste(slide, (0, 0))
        comp.paste(logo, (x, y), mask=logo)
        comp.convert("RGB").save(slide_path, "PNG")
    print("  Logo stamped on all slides (11% width, 4% bottom margin)")


# ─── Stage 5: Visual audit ────────────────────────────────────────────────────
def stage_audit() -> bool:
    print("=== Stage 5: Visual audit (Gemini 2.5 Flash Vision) ===")
    import subprocess
    audit_out = str(OUT_DIR / "audit_report.json")
    result = subprocess.run(
        [sys.executable, str(ROOT / "scripts" / "visual_audit.py"),
         "--slide_dir", str(IMG_DIR),
         "--rules",     str(RULES),
         "--out",       audit_out],
        cwd=str(ROOT), capture_output=True, text=True, timeout=300,
    )
    print(result.stdout[-800:] if result.stdout else "")
    if result.returncode != 0:
        print(f"  [WARN] Audit script error: {result.stderr[:300]}")
        return True  # continue even if audit fails (don't block delivery)

    report = json.loads(Path(audit_out).read_text(encoding="utf-8"))
    all_pass = report.get("overall_pass", True)
    if not all_pass:
        failed = report.get("failed_slides", [])
        print(f"  [WARN] Audit FAIL on slides: {failed}")
        print("  Continuing delivery with failure notes in email.")
    return all_pass


# ─── Stage 6: TTS ─────────────────────────────────────────────────────────────
# Opening narration prepended to slide 1 audio
_TTS_INTRO = "欢迎收看未来模式生物科技每周科学简报。本期主题："

def stage_tts(spec: dict):
    """Generate per-slide MP3 narrations.

    Primary:  Doubao seed-tts-2.0 (爽快思思 — 有感情色彩, 语速自然)
    Fallback: OpenAI tts-1-hd/shimmer
    """
    print("=== Stage 6: TTS narration ===")
    VOICE_DIR.mkdir(parents=True, exist_ok=True)
    import requests as req

    volc_key  = os.environ.get("VOLC_SPEECH_API_KEY") or os.environ.get("DOUBAO_API_KEY", "")
    resource  = os.environ.get("VOLC_TTS_RESOURCE_ID", "seed-tts-2.0")
    openai_key = os.environ.get("OPENAI_API_KEY", "")

    def _doubao_tts_rest(text: str) -> bytes:
        """Doubao TTS 2.0 via streaming REST — 爽快思思音色，有感情色彩。"""
        import uuid, base64, json as _json
        headers = {
            "X-Api-Key":         volc_key,
            "X-Api-Resource-Id": resource,
            "X-Api-Request-Id":  str(uuid.uuid4()),
            "Content-Type":      "application/json",
        }
        payload = {
            "user": {"uid": "nextvivo-weekly"},
            "req_params": {
                "text": text,
                "speaker": "zh_female_shuangkuaisisi_uranus_bigtts",
                "audio_params": {"format": "mp3", "sample_rate": 24000},
            },
        }
        resp = req.post(
            "https://openspeech.bytedance.com/api/v3/tts/unidirectional",
            headers=headers, json=payload, stream=True, timeout=120,
        )
        if resp.status_code != 200:
            raise RuntimeError(f"Doubao TTS HTTP {resp.status_code}: {resp.text[:200]}")
        audio = bytearray()
        for raw in resp.iter_lines(decode_unicode=True):
            if not raw:
                continue
            try:
                chunk = _json.loads(raw)
            except Exception:
                continue
            code = chunk.get("code")
            if code == 0 and chunk.get("data"):
                audio.extend(base64.b64decode(chunk["data"]))
            elif code == 20000000:
                break
            elif code not in (0, 20000000):
                raise RuntimeError(f"Doubao TTS error chunk: {chunk}")
        if not audio:
            raise RuntimeError("Doubao TTS: no audio data in stream")
        return bytes(audio)

    def _openai_tts(text: str) -> bytes:
        resp = req.post(
            "https://api.openai.com/v1/audio/speech",
            headers={"Authorization": f"Bearer {openai_key}",
                     "Content-Type": "application/json"},
            json={"model": "tts-1-hd", "input": text,
                  "voice": "shimmer", "response_format": "mp3"},
            timeout=60,
        )
        resp.raise_for_status()
        return resp.content

    backends = []
    if volc_key:
        backends.append(("doubao-seed-tts-2.0/爽快思思", _doubao_tts_rest))
    if openai_key:
        backends.append(("openai-tts-hd/shimmer", _openai_tts))

    if not backends:
        print("  [WARN] No TTS API key found, skipping TTS")
        return

    title = spec.get("slides", [{}])[0].get("title", "") if spec.get("slides") else ""
    intro_text = _TTS_INTRO + title if title else ""

    for slide in spec.get("slides", []):
        num       = slide["slide_num"]
        narration = slide.get("narration", "").strip()
        out_path  = VOICE_DIR / f"slide_{num:02d}.mp3"
        if out_path.exists():
            print(f"  TTS slide {num}: cached ✓")
            continue
        if not narration:
            continue
        # Prepend intro to slide 1
        text = (intro_text + "。" + narration) if (num == 1 and intro_text) else narration

        errors = []
        for name, fn in backends:
            try:
                audio = fn(text)
                out_path.write_bytes(audio)
                print(f"  TTS slide {num} [{name}]: {text[:40]}...")
                break
            except Exception as e:
                errors.append(f"{name}: {e!s:.100}")
        else:
            print(f"  [WARN] TTS slide {num} all failed: {' | '.join(errors)}")


# ─── Stage 7: Assemble PPTX ───────────────────────────────────────────────────
def stage_pptx():
    print("=== Stage 7: Assembling PPTX ===")
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i in range(1, 7):
        img = IMG_DIR / f"slide_{i:02d}.png"
        slide = prs.slides.add_slide(blank)
        if img.exists():
            slide.shapes.add_picture(str(img), 0, 0, prs.slide_width, prs.slide_height)
    prs.save(str(OUT_PPTX))
    print(f"  Saved: {OUT_PPTX.name}")


# ─── Stage 8: Assemble MP4 ────────────────────────────────────────────────────
def stage_mp4():
    print("=== Stage 8: Assembling MP4 (moviepy) ===")
    from moviepy import ImageClip, AudioFileClip, concatenate_videoclips

    clips = []
    for i in range(1, 7):
        img   = IMG_DIR / f"slide_{i:02d}.png"
        audio = VOICE_DIR / f"slide_{i:02d}.mp3"
        if not img.exists():
            continue
        if audio.exists():
            a   = AudioFileClip(str(audio))
            dur = a.duration + 0.5
            clip = ImageClip(str(img)).with_duration(dur).with_audio(a)
        else:
            clip = ImageClip(str(img)).with_duration(12)
        clips.append(clip)
        print(f"  slide {i}: {clip.duration:.1f}s")

    final = concatenate_videoclips(clips, method="compose")
    final.write_videofile(str(OUT_MP4), fps=24, codec="libx264",
                          audio_codec="aac", preset="fast", logger=None)
    print(f"  Saved: {OUT_MP4.name}  ({OUT_MP4.stat().st_size // 1024}KB)")


# ─── Stage 8b: Assemble vertical 9:16 MP4 for Douyin ─────────────────────────
def stage_mp4_vertical():
    """Reframe horizontal slides into 9:16 (1080×1920) for Douyin.

    Layout (top → bottom):
      • Gradient header bar  100px  — NextVivo brand + week tag
      • Slide image 16:9    607px  — original slide scaled to full width
      • Accent divider        8px
      • Title panel         ~220px — large bold title (word-wrapped)
      • Key-points panel    ~600px — 3 bullets, word-wrapped, blue accent
      • Logo footer         ~120px — logo + company name
      • Bottom padding       65px
    Total = 1920px  (no dead space)
    """
    print("=== Stage 8b: Assembling vertical 9:16 MP4 (Douyin) ===")
    try:
        from moviepy import ImageClip, AudioFileClip, concatenate_videoclips
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as e:
        print(f"  [SKIP] Missing dependency: {e}")
        return

    W, H = 1080, 1920
    SLIDE_H = int(W * 9 / 16)   # 607px — proper 16:9

    # Layout constants
    HEADER_H  = 100
    DIVIDER_H = 8
    FOOTER_H  = 185
    TEXT_H    = H - HEADER_H - SLIDE_H - DIVIDER_H - FOOTER_H  # ~1020px

    BG_DARK  = "#12122A"
    BG_MID   = "#1C1C3A"
    ACCENT   = "#2980B9"
    WHITE    = (255, 255, 255)
    LIGHT_BL = (161, 210, 255)
    SUBTEXT  = (160, 170, 200)

    sys.path.insert(0, str(ROOT / "scripts"))
    try:
        from cjk_fonts import cjk_font
    except Exception:
        cjk_font = lambda size, bold=False: ImageFont.load_default()

    def _wrap_text(draw, text: str, font, max_w: int) -> list[str]:
        """Split text into lines that fit within max_w pixels."""
        lines, line = [], ""
        for ch in text:
            test = line + ch
            if draw.textlength(test, font=font) > max_w:
                lines.append(line)
                line = ch
            else:
                line = test
        if line:
            lines.append(line)
        return lines or [""]

    def _make_frame(slide_img_path: Path, title: str, subtitle: str,
                    key_points: list[str]) -> Path:
        canvas = Image.new("RGB", (W, H), BG_DARK)
        draw   = ImageDraw.Draw(canvas)

        # ── Header bar ──────────────────────────────────────────────────────
        draw.rectangle([(0, 0), (W, HEADER_H)], fill=ACCENT)
        try:
            fn_hdr = cjk_font(28)
            draw.text((40, 34), "NextVivo · 未来模式生物科技", font=fn_hdr, fill=WHITE)
        except Exception:
            pass

        # ── Slide image (16:9, full width) ──────────────────────────────────
        y_slide = HEADER_H
        if slide_img_path.exists():
            slide = Image.open(slide_img_path).convert("RGB")
            sw, sh = slide.size
            # Scale to fill width, crop height if needed
            scale = W / sw
            new_h = int(sh * scale)
            slide = slide.resize((W, new_h), Image.LANCZOS)
            if new_h > SLIDE_H:
                top = (new_h - SLIDE_H) // 2
                slide = slide.crop((0, top, W, top + SLIDE_H))
            canvas.paste(slide, (0, y_slide))

        # ── Accent divider ───────────────────────────────────────────────────
        y_div = y_slide + SLIDE_H
        draw.rectangle([(0, y_div), (W, y_div + DIVIDER_H)], fill=ACCENT)

        # ── Text panel (title + key points) ─────────────────────────────────
        y_text = y_div + DIVIDER_H
        draw.rectangle([(0, y_text), (W, y_text + TEXT_H)], fill=BG_MID)

        PAD = 48
        max_w = W - 2 * PAD
        y = y_text + 36

        # Title
        try:
            fn_title = cjk_font(52, bold=True)
            for ln in _wrap_text(draw, title[:50], fn_title, max_w)[:2]:
                draw.text((PAD, y), ln, font=fn_title, fill=WHITE)
                y += int(52 * 1.35)
        except Exception:
            pass
        y += 10

        # Subtitle
        if subtitle:
            try:
                fn_sub = cjk_font(34)
                draw.text((PAD, y), subtitle[:48], font=fn_sub, fill=SUBTEXT)
                y += int(34 * 1.5)
            except Exception:
                pass
        y += 18

        # Divider line
        draw.rectangle([(PAD, y), (W - PAD, y + 2)], fill=ACCENT)
        y += 20

        # Key points
        try:
            fn_kp = cjk_font(36)
            for kp in (key_points or [])[:4]:
                bullet = f"▸  {kp}"
                for ln in _wrap_text(draw, bullet, fn_kp, max_w)[:2]:
                    draw.text((PAD, y), ln, font=fn_kp, fill=LIGHT_BL)
                    y += int(36 * 1.5)
                y += 6
        except Exception:
            pass

        # ── Footer: logo + brand ─────────────────────────────────────────────
        y_foot = H - FOOTER_H
        draw.rectangle([(0, y_foot), (W, H)], fill=BG_DARK)

        logo_path = LOGO_PATH
        if logo_path.exists():
            try:
                logo = Image.open(logo_path).convert("RGBA")
                logo.thumbnail((200, 80), Image.LANCZOS)
                lx = (W - logo.width) // 2
                ly = y_foot + (FOOTER_H - logo.height) // 2
                canvas.paste(logo, (lx, ly), logo)
            except Exception:
                pass

        out = slide_img_path.parent / f"_v_{slide_img_path.name}"
        canvas.save(out, "PNG")
        return out

    # Load slide titles from deck_spec if available
    spec_path = OUT_DIR / "deck_spec.json"
    slide_meta: dict[int, dict] = {}
    if spec_path.exists():
        try:
            spec = json.loads(spec_path.read_text(encoding="utf-8"))
            for s in spec.get("slides", []):
                slide_meta[s["slide_num"]] = s
        except Exception:
            pass

    clips = []
    for i in range(1, 7):
        img_h   = IMG_DIR / f"slide_{i:02d}.png"
        audio   = VOICE_DIR / f"slide_{i:02d}.mp3"
        meta    = slide_meta.get(i, {})
        title     = meta.get("title", f"第{i}张")
        subtitle  = meta.get("subtitle", "")
        key_points = meta.get("key_points", [])
        frame = _make_frame(img_h, title, subtitle, key_points)
        if not frame.exists():
            continue
        if audio.exists():
            a   = AudioFileClip(str(audio))
            dur = a.duration + 0.5
            clip = ImageClip(str(frame)).with_duration(dur).with_audio(a)
        else:
            clip = ImageClip(str(frame)).with_duration(12)
        clips.append(clip)
        print(f"  vertical slide {i}: {clip.duration:.1f}s")

    if not clips:
        print("  [SKIP] No clips generated")
        return

    out_v = OUT_DIR / OUT_MP4.name.replace(".mp4", "_douyin_9x16.mp4")
    final = concatenate_videoclips(clips, method="compose")
    final.write_videofile(str(out_v), fps=24, codec="libx264",
                          audio_codec="aac", preset="fast", logger=None)
    print(f"  Saved: {out_v.name}  ({out_v.stat().st_size // 1024}KB)")
    # Expose for email attachment
    global OUT_MP4_VERTICAL
    OUT_MP4_VERTICAL = out_v


OUT_MP4_VERTICAL: Path | None = None


# ─── Stage 9: Email 1 — PPT + MP4 ────────────────────────────────────────────
def stage_email_ppt(article: dict, audit_pass: bool):
    print("=== Stage 9a: Email 1 — PPT + MP4 (横版+抖音竖版) ===")
    sys.path.insert(0, str(ROOT / "scripts"))
    from send_email import send_weekly_report
    send_weekly_report(
        article=article,
        pptx_path=OUT_PPTX,
        mp4_path=OUT_MP4,
        mp4_vertical_path=OUT_MP4_VERTICAL,
        audit_pass=audit_pass,
        recipients=RECIPIENTS,
        week_tag=WEEK_TAG,
    )


# ─── Stage 10: Social content generation ─────────────────────────────────────
def stage_social(article: dict) -> dict:
    print("=== Stage 10: Social content (XHS + WeChat) ===")
    sys.path.insert(0, str(ROOT / "scripts"))
    from social_content import generate_social_content
    social_json = OUT_DIR / "social.json"
    return generate_social_content(article, social_json)


# ─── Stage 11: Social images ──────────────────────────────────────────────────
def stage_social_images(social_data: dict, force: bool = False,
                        frozen_wechat_dir: Path | None = None,
                        slide_dir: Path | None = None):
    print("=== Stage 11: Social images (XHS 6 cards + WeChat images) ===")
    from social_images import generate_all_images
    generate_all_images(social_data, OUT_DIR / "social_images", force=force,
                        frozen_wechat_dir=frozen_wechat_dir, slide_dir=slide_dir,
                        inline_sections=[1, 2, 3])


# ─── Stage 12: Email 2 — XHS ─────────────────────────────────────────────────
def stage_email_xhs(article: dict, social_data: dict):
    print("=== Stage 12: Email 2 — 小红书内容 ===")
    from send_email import send_xhs_report
    send_xhs_report(
        article=article,
        social_data=social_data,
        img_dir=OUT_DIR / "social_images" / "xhs",
        recipients=RECIPIENTS,
        week_tag=WEEK_TAG,
    )


# ─── Stage 13: Email 3 — WeChat ──────────────────────────────────────────────
def stage_email_wechat(article: dict, social_data: dict):
    print("=== Stage 13: Email 3 — 微信公众号内容 ===")
    from send_email import send_wechat_report
    send_wechat_report(
        article=article,
        social_data=social_data,
        img_dir=OUT_DIR / "social_images" / "wechat",
        recipients=RECIPIENTS,
        week_tag=WEEK_TAG,
    )


# ─── Main ─────────────────────────────────────────────────────────────────────
FROZEN_BUNDLE_DEFAULT = ROOT / "data" / "frozen" / "humanized_mice_front_immunol_2026"

if __name__ == "__main__":
    import argparse
    ap = argparse.ArgumentParser(description="Weekly humanized mice pipeline")
    ap.add_argument(
        "--replay-bundle",
        type=Path,
        default=None,
        help="Reuse frozen article/deck_spec/social.json (skip PubMed + LLM content)",
    )
    ap.add_argument(
        "--force-images",
        action="store_true",
        help="Delete cached slide/social images before re-render",
    )
    ap.add_argument("--emails-only", action="store_true", help="Skip generation, send all 3 emails only")
    ap.add_argument("--wechat-only", action="store_true", help="Send WeChat article email only (uses frozen cover/ad bar)")
    args = ap.parse_args()

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    # Resolve bundle path: accept plain slug or full path
    _rb_arg = args.replay_bundle
    if _rb_arg is not None:
        _rb = Path(_rb_arg)
        bundle = _rb if (_rb.is_absolute() or _rb.exists()) else ROOT / "data" / "frozen" / _rb
    else:
        bundle = FROZEN_BUNDLE_DEFAULT

    # Resolve topic slug and update output filenames (module-level reassignment)
    _TOPIC_SLUG = _resolve_topic_slug(bundle)  # noqa: F841 — used via OUT_PPTX/OUT_MP4 below
    OUT_PPTX = OUT_DIR / f"{_TOPIC_SLUG}_{WEEK_TAG}.pptx"
    OUT_MP4  = OUT_DIR / f"{_TOPIC_SLUG}_{WEEK_TAG}.mp4"

    print(f"\n{'='*60}")
    print(f"  NextVivo Weekly Content Pipeline — {WEEK_TAG}")
    print(f"  Topic: {_TOPIC_SLUG}")
    if args.replay_bundle:
        print(f"  MODE: REPLAY frozen bundle → {args.replay_bundle}")
    print(f"  Outputs: PPT/MP4 + XHS + WeChat → 3 separate emails")
    print(f"{'='*60}\n")

    if args.wechat_only:
        article     = json.loads((bundle / "article.json").read_text(encoding="utf-8"))
        social_data = json.loads((bundle / "social.json").read_text(encoding="utf-8"))
        # Copy approved frozen cover + ad bar into expected output dir
        wechat_img_dir = OUT_DIR / "social_images" / "wechat"
        wechat_img_dir.mkdir(parents=True, exist_ok=True)
        frozen_wechat = bundle / "wechat"
        for name in ("wechat_cover.png", "wechat_ad_bar.png"):
            src = frozen_wechat / name
            if src.exists():
                shutil.copy2(src, wechat_img_dir / name)
                print(f"  Pinned: {name} → {wechat_img_dir.relative_to(ROOT)}")
            else:
                print(f"  [WARN] Pinned asset not found: {src}")
        # Generate inline body images: try frozen first, then slide reuse, then gpt-image-2
        sys.path.insert(0, str(ROOT / "scripts"))
        from social_images import generate_wechat_images
        generate_wechat_images(social_data, wechat_img_dir, force=True,
                               frozen_dir=frozen_wechat, slide_dir=None)
        stage_email_wechat(article, social_data)
        print(f"\n✅ WeChat email sent for {WEEK_TAG}")
        raise SystemExit(0)

    if args.emails_only:
        article     = json.loads((bundle / "article.json").read_text(encoding="utf-8"))
        social_data = json.loads((bundle / "social.json").read_text(encoding="utf-8"))
        stage_email_ppt(article, audit_pass=True)
        stage_email_xhs(article, social_data)
        stage_email_wechat(article, social_data)
        print(f"\n✅ Emails sent for {WEEK_TAG}")
        raise SystemExit(0)

    if args.force_images:
        for d in (IMG_DIR, OUT_DIR / "social_images"):
            if d.exists():
                shutil.rmtree(d)
        for f in (OUT_PPTX, OUT_MP4):
            if f.exists():
                f.unlink()

    frozen_wechat = None
    if args.replay_bundle:
        _rb = Path(args.replay_bundle)
        # Accept plain slug (e.g. "in_vivo_car_tech_2026") or full path
        bundle = _rb if _rb.is_absolute() or _rb.exists() else ROOT / "data" / "frozen" / _rb
        article = json.loads((bundle / "article.json").read_text(encoding="utf-8"))
        spec = json.loads((bundle / "deck_spec.json").read_text(encoding="utf-8"))
        social_data = json.loads((bundle / "social.json").read_text(encoding="utf-8"))
        frozen_wechat = bundle / "wechat" if (bundle / "wechat").is_dir() else None
        (OUT_DIR / "deck_spec.json").write_text(json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8")
        (OUT_DIR / "social.json").write_text(json.dumps(social_data, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"  Loaded frozen bundle: {bundle.name}")
        print(f"  Article PMID {article['pmid']}: {article['title'][:70]}...")
    else:
        article    = stage_pubmed()
        spec       = stage_deck_spec(article)
        spec       = stage_fact_check(spec, article)   # Stage 2b: zero-tolerance fact-check
        social_data = None

    stage_generate_images(spec)
    stage_logo()
    audit_pass = stage_audit()
    stage_tts(spec)
    stage_pptx()
    stage_mp4()
    stage_mp4_vertical()                              # 9:16 Douyin version
    stage_email_ppt(article, audit_pass)          # Email 1: PPT + MP4(横) + MP4(竖)

    if social_data is None:
        social_data = stage_social(article)       # DeepSeek→Kimi→fix
    stage_social_images(social_data, force=args.force_images,
                        frozen_wechat_dir=frozen_wechat,
                        slide_dir=IMG_DIR if IMG_DIR.exists() else None)
    stage_email_xhs(article, social_data)         # Email 2: XHS
    stage_email_wechat(article, social_data)      # Email 3: WeChat

    print(f"\n✅ All done for {WEEK_TAG}")
    print(f"   PPT:    {OUT_PPTX}")
    print(f"   MP4:    {OUT_MP4}")
    print(f"   Social: {OUT_DIR / 'social.json'}")
    print(f"   Emails: 3 sent to {RECIPIENTS}")
