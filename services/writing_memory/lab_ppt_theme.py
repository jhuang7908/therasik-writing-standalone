"""InSynBio Lab — default Project Report PPT template (scientific blue house style).

This module centralizes the visual theme for generated decks so that restyling
(brand colors, fonts, logo, layout) is a single-file change and never requires
touching the report-extraction logic in ``app.py``.

Design goals:
- Free + offline (pure python-pptx, no external services).
- Consistent "house style": branded cover, section dividers, accent header bar,
  branded footer with slide numbers, professional zebra tables, captioned figures.
- Defensive: every builder tolerates empty/missing content without raising.
"""

from __future__ import annotations

import base64
import io
import re
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── InSynBio scientific-blue palette ─────────────────────────────────────────
PRIMARY = RGBColor(0x1D, 0x4E, 0x89)   # deep scientific blue (headers, cover)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)    # lighter teal-blue (bars, highlights)
DARK = RGBColor(0x1A, 0x20, 0x2A)      # body text
GREY = RGBColor(0x5A, 0x63, 0x70)      # secondary / captions
LIGHT = RGBColor(0xEE, 0xF3, 0xF8)     # zebra row / panel background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEAD_FONT = "Calibri"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class LabReportDeck:
    """Builds a themed 16:9 deck. Call builders in order, then ``save()``."""

    def __init__(self, report_id: str = "", brand: str = "InSynBio Lab", generated: str = "") -> None:
        self.report_id = report_id or ""
        self.brand = brand
        self.generated = (generated or "")[:10]
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]
        self._n = 0  # content slide counter (cover excluded)

    # ── low-level helpers ────────────────────────────────────────────────
    def _rect(self, slide, x, y, w, h, color):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _run(self, paragraph, text, *, size, color, bold=False, font=BODY_FONT):
        r = paragraph.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
        return r

    def _content_slide(self, head: str):
        s = self.prs.slides.add_slide(self._blank)
        self._n += 1
        # thin accent bar at very top
        self._rect(s, Emu(0), Emu(0), SLIDE_W, Inches(0.16), ACCENT)
        # title
        box = s.shapes.add_textbox(Inches(0.55), Inches(0.38), Inches(12.2), Inches(0.75))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        self._run(p, (head or "")[:90], size=27, color=PRIMARY, bold=True, font=HEAD_FONT)
        # underline rule beneath title
        self._rect(s, Inches(0.6), Inches(1.18), Inches(3.2), Pt(2.4), PRIMARY)
        self._footer(s)
        return s

    def _footer(self, slide):
        box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(10.5), Inches(0.3))
        p = box.text_frame.paragraphs[0]
        bits = [b for b in [self.brand, self.report_id and f"Report {self.report_id}", self.generated] if b]
        self._run(p, "  ·  ".join(bits), size=8, color=GREY)
        # slide number (right)
        nbox = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.3))
        np = nbox.text_frame.paragraphs[0]
        np.alignment = PP_ALIGN.RIGHT
        self._run(np, str(self._n), size=8, color=GREY)

    # ── public builders ──────────────────────────────────────────────────
    def cover(self, title: str, subtitle: str = "", qc: str = "", tagline: str = "") -> None:
        s = self.prs.slides.add_slide(self._blank)
        # full-bleed primary background
        self._rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, PRIMARY)
        # accent stripe
        self._rect(s, Emu(0), Inches(2.35), SLIDE_W, Inches(0.08), ACCENT)
        # white title card
        self._rect(s, Emu(0), Inches(2.55), SLIDE_W, Inches(2.7), WHITE)

        tb = s.shapes.add_textbox(Inches(0.9), Inches(2.85), Inches(11.5), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        self._run(tf.paragraphs[0], (title or "Project Report")[:140], size=38, color=PRIMARY, bold=True, font=HEAD_FONT)

        if subtitle:
            sp = tf.add_paragraph()
            self._run(sp, subtitle[:160], size=15, color=GREY)

        # QC badge
        if qc:
            badge = self._rect(s, Inches(0.9), Inches(4.55), Inches(2.4), Inches(0.5), ACCENT)
            bp = badge.text_frame.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            badge.text_frame.word_wrap = True
            self._run(bp, f"QC: {qc}", size=14, color=WHITE, bold=True)

        # top-left brand
        brand_box = s.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.6))
        self._run(brand_box.text_frame.paragraphs[0], self.brand, size=20, color=WHITE, bold=True, font=HEAD_FONT)
        if tagline:
            tlb = s.shapes.add_textbox(Inches(0.9), Inches(1.25), Inches(11.5), Inches(0.5))
            self._run(tlb.text_frame.paragraphs[0], tagline[:120], size=12, color=LIGHT)

        # bottom meta
        meta = s.shapes.add_textbox(Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5))
        bits = [b for b in [self.generated and f"Generated {self.generated}",
                            self.report_id and f"Archive {self.report_id}"] if b]
        self._run(meta.text_frame.paragraphs[0], "   ·   ".join(bits), size=11, color=LIGHT)

    def section_divider(self, text: str) -> None:
        s = self.prs.slides.add_slide(self._blank)
        self._rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, PRIMARY)
        self._rect(s, Inches(0.9), Inches(3.35), Inches(2.6), Inches(0.09), ACCENT)
        tb = s.shapes.add_textbox(Inches(0.9), Inches(3.55), Inches(11.5), Inches(1.2))
        tb.text_frame.word_wrap = True
        self._run(tb.text_frame.paragraphs[0], (text or "")[:80], size=34, color=WHITE, bold=True, font=HEAD_FONT)

    def bullets(self, head: str, body: str, *, max_chars: int = 1600, max_lines: int = 11) -> None:
        body = (body or "").strip()
        if not body:
            return
        s = self._content_slide(head)
        box = s.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(11.95), Inches(5.3))
        tf = box.text_frame
        tf.word_wrap = True
        lines = re.split(r"\n+|(?<=[.;。！？])\s+", body[:max_chars])
        lines = [x.strip() for x in lines if x.strip()][:max_lines]
        first = True
        for line in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(7)
            self._run(p, "•  ", size=15, color=ACCENT, bold=True)
            self._run(p, line, size=15, color=DARK)

    def table(self, caption: str, rows: list[list[str]]) -> None:
        rows = [r for r in (rows or []) if r]
        if not rows:
            return
        s = self._content_slide(caption or "Data Table")
        n_rows = min(len(rows), 11)
        n_cols = min(max(len(r) for r in rows), 6)
        gfx = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.4) * n_rows)
        table = gfx.table
        # disable python-pptx default banded style so our fills show cleanly
        for i in range(n_rows):
            row = rows[i]
            table.rows[i].height = Inches(0.42)
            for j in range(n_cols):
                cell = table.cell(i, j)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_top = Pt(2)
                cell.margin_bottom = Pt(2)
                val = row[j] if j < len(row) else ""
                cell.text = str(val)[:120]
                cell.fill.solid()
                if i == 0:
                    cell.fill.fore_color.rgb = PRIMARY
                else:
                    cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
                para = cell.text_frame.paragraphs[0]
                para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                for r in para.runs:
                    r.font.size = Pt(12 if i == 0 else 11)
                    r.font.bold = i == 0
                    r.font.name = BODY_FONT
                    r.font.color.rgb = WHITE if i == 0 else DARK

    def figure(self, caption: str, src: str, *, note: str = "") -> None:
        s = self._content_slide(caption or "Figure")
        src = src or ""
        placed = False
        if src.startswith("data:image"):
            try:
                raw = base64.b64decode(src.split(",", 1)[1])
                pic = s.shapes.add_picture(io.BytesIO(raw), Inches(0.8), Inches(1.5))
                # scale to fit a 11.7 x 5.0 in box preserving aspect
                max_w, max_h = Inches(11.7), Inches(5.0)
                scale = min(max_w / pic.width, max_h / pic.height, 1.0)
                pic.width = int(pic.width * scale)
                pic.height = int(pic.height * scale)
                pic.left = int((SLIDE_W - pic.width) / 2)
                pic.top = Inches(1.55)
                placed = True
            except Exception:
                placed = False
        if not placed:
            panel = self._rect(s, Inches(1.4), Inches(2.7), Inches(10.5), Inches(1.6), LIGHT)
            tf = panel.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            msg = note or (
                "No embedded chart image in this archived report. "
                "Generate the chart in Experimental Data (upload CSV → Draw chart) "
                "to embed Python figures here."
            )
            self._run(tf.paragraphs[0], msg, size=15, color=GREY)

    def save(self) -> bytes:
        buf = io.BytesIO()
        self.prs.save(buf)
        buf.seek(0)
        return buf.read()

    @property
    def n_slides(self) -> int:
        return len(self.prs.slides._sldIdLst)  # noqa: SLF001
